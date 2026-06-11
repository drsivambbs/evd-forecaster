"""Standalone PowerPoint report builder for the EVD Forecaster.

Design goals:
  * No Streamlit or Plotly import — charts are drawn with matplotlib so image
    export never depends on kaleido/Chrome (which hangs on this box with
    plotly 6 + kaleido 0.2.1). The module is therefore importable and testable
    on its own.
  * One slide per step: inputs table (left) | chart (right) | interpretation
    (bottom). Plus a title slide, an executive-summary slide and a
    methods/references slide.

Entry point:  build_pptx(ss, scenario_name) -> bytes
`ss` is any Mapping with the app's session_state values (st.session_state works
directly). Every renderer degrades gracefully: a missing/!broken chart becomes
a placeholder note rather than an exception.
"""
from __future__ import annotations

from io import BytesIO
from datetime import datetime

import numpy as np
import pandas as pd

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.dates as mdates

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --------------------------------------------------------------------------
# Brand / palette
# --------------------------------------------------------------------------
BRAND = RGBColor(0x1F, 0x4E, 0x79)
BRAND_HEX = "#1f4e79"
GREY = RGBColor(0x5B, 0x65, 0x73)
LIGHT = RGBColor(0xF1, 0xF4, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xE3, 0xE7, 0xED)

COLOURS = {"confirmed": "#4682B4", "suspected": "#FF8C00", "deaths": "#B22222"}
SCEN_COLOURS = {"S1": "#C0392B", "S2": "#E07B39", "S3": "#1E8449"}
RESP = {"S1": "Delayed response", "S2": "Moderate response",
        "S3": "Strong response"}

plt.rcParams.update({
    "font.family": "DejaVu Sans",
    "font.size": 11,
    "axes.edgecolor": "#cfd6df",
    "axes.titlecolor": BRAND_HEX,
    "axes.titlesize": 13,
    "axes.labelcolor": "#333333",
    "axes.grid": True,
    "grid.color": "#eef1f5",
    "figure.facecolor": "white",
    "axes.facecolor": "white",
})


# ==========================================================================
# Matplotlib chart renderers  ->  PNG bytes (or None on failure)
# ==========================================================================
def _save(fig) -> bytes:
    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor="white")
    plt.close(fig)
    return buf.getvalue()


def _date_axis(ax):
    ax.xaxis.set_major_formatter(mdates.DateFormatter("%d %b %y"))
    for lbl in ax.get_xticklabels():
        lbl.set_rotation(0)
        lbl.set_fontsize(9)


def _col(series, *candidates):
    for c in candidates:
        if c in series.columns:
            return c
    return None


def render_daily(series) -> bytes | None:
    try:
        s = series.copy()
        s["date"] = pd.to_datetime(s["date"])
        fig, ax = plt.subplots(figsize=(7.2, 4.2))
        plots = [
            ("Actual confirmed", _col(s, "new_confirmed_raw", "new_confirmed"),
             COLOURS["confirmed"], "-"),
            ("Actual suspected", _col(s, "new_suspected_raw", "new_suspected"),
             COLOURS["suspected"], "-"),
            ("Actual deaths", _col(s, "new_deaths_raw", "new_deaths"),
             COLOURS["deaths"], "-"),
            ("Estimated confirmed", _col(s, "new_cfr_estimated"),
             "#6a1b9a", "--"),
            ("Estimated suspected", _col(s, "new_suspected_estimated"),
             "#b8860b", "--"),
        ]
        for label, col, colour, style in plots:
            if col is None:
                continue
            if float(np.nan_to_num(s[col]).sum()) == 0:
                continue
            ax.plot(s["date"], s[col], style, color=colour, linewidth=2.0,
                    label=label)
        ax.set_title("Daily new cases", loc="left")
        ax.set_ylabel("New cases per day")
        ax.legend(fontsize=9, frameon=False, ncol=2)
        _date_axis(ax)
        return _save(fig)
    except Exception:
        return None


def render_cumulative(series, snaps) -> bytes | None:
    try:
        s = series.copy()
        s["date"] = pd.to_datetime(s["date"])
        fig, ax = plt.subplots(figsize=(7.2, 4.2))
        for key, label in [("cumulative_confirmed", "Confirmed"),
                           ("cumulative_suspected", "Suspected"),
                           ("cumulative_deaths", "Deaths")]:
            if key in s.columns:
                ax.plot(s["date"], s[key], "-",
                        color=COLOURS[label.lower()], linewidth=2.0,
                        label=label)
        if snaps is not None and len(snaps) > 0:
            snp = snaps.dropna(subset=["cumulative_confirmed",
                                       "cumulative_suspected",
                                       "cumulative_deaths"]).copy()
            snp["date"] = pd.to_datetime(snp["date"])
            snp = snp.sort_values("date")
            for key, label in [("cumulative_confirmed", "confirmed"),
                               ("cumulative_suspected", "suspected"),
                               ("cumulative_deaths", "deaths")]:
                ax.scatter(snp["date"], snp[key], s=22,
                           color=COLOURS[label], edgecolor="white",
                           zorder=5)
        ax.set_title("Cumulative cases (interpolated + snapshots)", loc="left")
        ax.set_ylabel("Cumulative count")
        ax.legend(fontsize=9, frameon=False)
        _date_axis(ax)
        return _save(fig)
    except Exception:
        return None


def render_rt(rt_df, si_mean) -> bytes | None:
    try:
        prim = rt_df[rt_df["si_mean_used"] == si_mean].copy()
        prim["date"] = pd.to_datetime(prim["date"])
        fig, ax = plt.subplots(figsize=(7.2, 4.2))
        ax.fill_between(prim["date"], prim["rt_lower"], prim["rt_upper"],
                        color="#B22222", alpha=0.18, label="95% CrI")
        ax.plot(prim["date"], prim["rt_mean"], "-", color="#8B0000",
                linewidth=2.4, marker="o", markersize=4, label="R_t mean")
        sens = rt_df[rt_df["si_mean_used"] != si_mean]
        if len(sens) > 0:
            sm = sens["si_mean_used"].iloc[0]
            sens = sens.copy()
            sens["date"] = pd.to_datetime(sens["date"])
            ax.plot(sens["date"], sens["rt_mean"], "--", color="#2E8B8B",
                    linewidth=1.8, label=f"Sensitivity (SI={sm:.1f} d)")
        ax.axhline(1.0, color="black", linestyle=":", linewidth=1.3,
                   label="R_t = 1")
        ax.set_title(f"Instantaneous R_t (SI mean = {si_mean:.1f} d)",
                     loc="left")
        ax.set_ylabel("Instantaneous R_t")
        ax.legend(fontsize=9, frameon=False)
        _date_axis(ax)
        return _save(fig)
    except Exception:
        return None


def _mlu(metric):
    if isinstance(metric, dict):
        return metric.get("median"), metric.get("lower"), metric.get("upper")
    return metric, None, None


def render_forecast(scenarios, dates, baselines, labels,
                    y_log=True) -> bytes | None:
    try:
        order = [s for s in ["S1", "S2", "S3"] if s in scenarios]
        dts = pd.to_datetime(list(dates))
        fig, axes = plt.subplots(1, len(order), figsize=(9.6, 4.2),
                                 sharey=True)
        if len(order) == 1:
            axes = [axes]
        metrics = [("cum_confirmed", "Confirmed", COLOURS["confirmed"]),
                   ("cum_suspected", "Suspected", COLOURS["suspected"]),
                   ("cum_deaths", "Deaths", COLOURS["deaths"])]
        for ax, name in zip(axes, order):
            data = scenarios[name]
            for key, label, colour in metrics:
                med, lo, hi = _mlu(data[key])
                if med is None:
                    continue
                if lo is not None and hi is not None:
                    ax.fill_between(dts, lo, hi, color=colour, alpha=0.16)
                ax.plot(dts, med, "-", color=colour, linewidth=2.0,
                        label=label)
            if y_log:
                ax.set_yscale("log")
            ax.set_title(labels.get(name, name), fontsize=11)
            ax.xaxis.set_major_formatter(mdates.DateFormatter("%b %y"))
            for lbl in ax.get_xticklabels():
                lbl.set_fontsize(8)
        axes[0].set_ylabel("Cumulative count" + (" (log)" if y_log else ""))
        axes[-1].legend(fontsize=9, frameon=False)
        fig.suptitle("Forecast — cumulative cases by response scenario",
                     x=0.01, ha="left", color=BRAND_HEX, fontsize=13)
        fig.tight_layout(rect=[0, 0, 1, 0.95])
        return _save(fig)
    except Exception:
        return None


def render_eoo(days, probs, plc, threshold) -> bytes | None:
    try:
        fig, ax = plt.subplots(figsize=(7.2, 4.2))
        days = np.asarray(days)
        for name, last in plc.items():
            x = [pd.Timestamp(last) + pd.Timedelta(days=int(d)) for d in days]
            ax.plot(x, probs, "-", color=SCEN_COLOURS.get(name, "#333"),
                    linewidth=2.4,
                    label=f"{RESP.get(name, name)} — last case "
                          f"{pd.Timestamp(last).strftime('%d %b %Y')}")
        ax.axhline(threshold, color="#7a7a7a", linestyle=":", linewidth=1.3)
        ax.text(0.995, threshold + 0.01,
                f"{int(threshold*100)}% declaration threshold",
                transform=ax.get_yaxis_transform(), ha="right", fontsize=8,
                color="#444")
        ax.set_ylim(0, 1.05)
        ax.set_title("End-of-outbreak probability per scenario", loc="left")
        ax.set_ylabel("P(outbreak extinct)")
        ax.legend(fontsize=8, frameon=False)
        _date_axis(ax)
        return _save(fig)
    except Exception:
        return None


# ==========================================================================
# Per-step input tables  &  interpretation text  (read from ss)
# ==========================================================================
def _g(ss, key, default=None):
    try:
        v = ss.get(key, default)
    except AttributeError:
        v = default
    return default if v is None else v


def step1_inputs(ss):
    return [("Scenario label", str(_g(ss, "scenario_name", "—")))]


def step2_inputs(ss):
    return [
        ("SI mean / SD (days)",
         f"{_g(ss, 'rt_si_mean', 15.3)} / {_g(ss, 'rt_si_sd', 9.3)}"),
        ("Sliding window (days)", f"{_g(ss, 'window_val', 7)}"),
        ("Gamma prior (shape, rate)",
         f"{_g(ss, 'shape_prior_val', 1.0)}, {_g(ss, 'rate_prior_val', 0.2)}"),
        ("R_t selected for forecast",
         f"{_g(ss, 'selected_rt', '—')}  ({_g(ss, 'selected_rt_basis', '—')})"),
    ]


def step3_inputs(ss):
    return [
        ("CFR (fraction)", f"{_g(ss, 'fc_cfr', 0.34)}"),
        ("Onset-to-death lag (days)", f"{_g(ss, 'lag_val', 10)}"),
        ("Forecast horizon (days)", f"{_g(ss, 'forecast_horizon', 365)}"),
        ("S1 target / days", f"{_g(ss, 'S1_target', 0.9)} / {_g(ss, 'S1_days', 180)}"),
        ("S2 target / days", f"{_g(ss, 'S2_target', 0.9)} / {_g(ss, 'S2_days', 90)}"),
        ("S3 target / days", f"{_g(ss, 'S3_target', 0.6)} / {_g(ss, 'S3_days', 30)}"),
    ]


def step4_inputs(ss):
    return [
        ("Simulations", f"{_g(ss, 'eoo_n_sim', 1000):,}"),
        ("Days range", f"{_g(ss, 'eoo_max_days', 180)}"),
        ("Declaration threshold", f"{_g(ss, 'eoo_threshold', 0.95):.0%}"),
        ("Offspring dispersion k", f"{_g(ss, 'eoo_k_disp', 0.18)}"),
    ]


def step1_interp(series, ss):
    d = pd.to_datetime(series["date"])
    def _sum(col):
        rc = f"{col}_raw"
        src = series[rc] if rc in series.columns else series.get(col)
        return float(src.sum()) if src is not None else 0.0
    return (
        f"Daily incidence spans {len(series)} days "
        f"({d.min():%d %b %Y} to {d.max():%d %b %Y}). Totals — confirmed "
        f"{_sum('new_confirmed'):,.0f}, suspected {_sum('new_suspected'):,.0f}, "
        f"deaths {_sum('new_deaths'):,.0f}. Counts between sparse cumulative "
        f"snapshots are linearly interpolated, then first-differenced.")


def step2_interp(rt_df, si, ss):
    prim = rt_df[rt_df["si_mean_used"] == si].dropna(subset=["rt_mean"])
    if prim.empty:
        return "R_t estimates computed (Cori et al. 2013 sliding window)."
    last = prim.iloc[-1]
    rt = float(last["rt_mean"])
    phase = "growing (R_t > 1)" if rt > 1 else "declining (R_t < 1)"
    return (
        f"Latest R_t = {rt:.2f} (95% CrI {float(last['rt_lower']):.2f}–"
        f"{float(last['rt_upper']):.2f}) on "
        f"{pd.to_datetime(last['date']):%d %b %Y} — the outbreak is {phase}. "
        f"Estimated with the Cori et al. (2013) sliding-window method on a "
        f"Gamma serial interval (mean {si:.1f} d).")


def step3_interp(scenarios, dates, ss):
    bits = []
    for name in [s for s in ["S1", "S2", "S3"] if s in scenarios]:
        med, _, _ = _mlu(scenarios[name]["cum_confirmed"])
        bits.append(f"{RESP.get(name, name)}: {float(med[-1]):,.0f} confirmed")
    return (
        "Renewal-equation forecast (Nouvellet et al. 2018). Cumulative confirmed "
        f"at horizon end — " + "; ".join(bits) + ". Shaded bands are 90% "
        "posterior predictive intervals from sampling the R_t starting value.")


def step4_interp(days, probs, plc, thr):
    probs = np.asarray(probs)
    days = np.asarray(days)
    cross = np.where(probs >= thr)[0]
    if len(cross) == 0:
        return (f"P(extinct) does not reach the {thr:.0%} threshold within "
                f"{int(days[-1])} days. Extend the days range or revisit the "
                "Step 3 scenarios.")
    cd = int(days[cross[0]])
    lines = []
    for name, last in plc.items():
        dt = (pd.Timestamp(last) + pd.Timedelta(days=cd)).strftime("%d %b %Y")
        lines.append(f"{RESP.get(name, name)}: {dt}")
    return (
        f"P(extinct) crosses {thr:.0%} at day {cd} after the projected last "
        f"case (Nishiura / Lloyd-Smith offspring-tree simulation). Declaration "
        f"dates — " + "; ".join(lines) + ".")


# ==========================================================================
# PPTX layout helpers
# ==========================================================================
SW, SH = Inches(13.333), Inches(7.5)


def _header(slide, title, subtitle=""):
    bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(0.95))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BRAND
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.4)
    tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = subtitle
        r2.font.size = Pt(12); r2.font.color.rgb = RGBColor(0xD6, 0xE0, 0xEE)


def _footer(slide, page):
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(7.05),
                                  Inches(12.5), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = ("EVD Forecaster  ·  Cori 2013 · Nouvellet 2018 · "
              "Lloyd-Smith 2005   |   " + str(page))
    r.font.size = Pt(8); r.font.color.rgb = GREY


def _table(slide, rows, x, y, w):
    n = len(rows) + 1
    h = Inches(min(0.34 * n + 0.1, 4.7))
    gt = slide.shapes.add_table(n, 2, x, y, w, h).table
    gt.columns[0].width = Inches(2.5)
    gt.columns[1].width = w - Inches(2.5)
    for j, head in enumerate(["Parameter", "Value"]):
        c = gt.cell(0, j)
        c.text = head
        c.fill.solid(); c.fill.fore_color.rgb = BRAND
        pr = c.text_frame.paragraphs[0].runs[0]
        pr.font.size = Pt(11); pr.font.bold = True; pr.font.color.rgb = WHITE
    for i, (k, v) in enumerate(rows, start=1):
        for j, val in enumerate([k, str(v)]):
            c = gt.cell(i, j)
            c.text = val
            c.fill.solid()
            c.fill.fore_color.rgb = LIGHT if i % 2 else WHITE
            run = c.text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(10)
            if j == 0:
                run.font.color.rgb = GREY
    return gt


def _interp_box(slide, text, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "Interpretation"
    r.font.bold = True; r.font.size = Pt(12); r.font.color.rgb = BRAND
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = text
    r2.font.size = Pt(11); r2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def _placeholder(slide, x, y, w, h, msg):
    tb = slide.shapes.add_textbox(x, y, w, h)
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = msg
    r.font.size = Pt(12); r.font.italic = True; r.font.color.rgb = GREY


def _step_slide(prs, title, input_rows, png, interp, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, title)
    _table(slide, input_rows, Inches(0.45), Inches(1.2), Inches(4.5))
    if png:
        slide.shapes.add_picture(BytesIO(png), Inches(5.2), Inches(1.15),
                                 width=Inches(7.7))
    else:
        _placeholder(slide, Inches(5.2), Inches(2.6), Inches(7.7),
                     Inches(1.0), "Chart unavailable for this run.")
    _interp_box(slide, interp, Inches(0.45), Inches(5.85), Inches(12.4),
                Inches(1.1))
    _footer(slide, page)


def _title_slide(prs, scenario_name, now):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(1, 0, Inches(2.4), SW, Inches(2.7))
    bar.fill.solid(); bar.fill.fore_color.rgb = BRAND
    bar.line.fill.background()
    tf = bar.text_frame; tf.margin_left = Inches(0.7)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "EVD Forecaster"
    r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = scenario_name
    r2.font.size = Pt(20); r2.font.color.rgb = RGBColor(0xD6, 0xE0, 0xEE)
    sub = slide.shapes.add_textbox(Inches(0.7), Inches(5.3), Inches(12), Inches(1))
    sp = sub.text_frame.paragraphs[0]
    sr = sp.add_run()
    sr.text = (f"Renewal-equation outbreak forecast report  ·  generated {now}")
    sr.font.size = Pt(13); sr.font.color.rgb = GREY


def _summary_slide(prs, ss, scenario_name):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "Executive summary", scenario_name)
    rows = []
    series = _g(ss, "result_series")
    if series is not None and len(series) > 0:
        rows.append(("Step 1 — Daily incidence",
                     f"{len(series)} days built"))
    rt_df = _g(ss, "rt_df")
    si = float(_g(ss, "rt_si_mean", 15.3))
    if rt_df is not None and len(rt_df) > 0:
        prim = rt_df[rt_df["si_mean_used"] == si].dropna(subset=["rt_mean"])
        if not prim.empty:
            rt = float(prim.iloc[-1]["rt_mean"])
            rows.append(("Step 2 — Latest R_t",
                         f"{rt:.2f} ({'growing' if rt > 1 else 'declining'})"))
    scen = _g(ss, "fc_scenarios")
    if scen:
        bits = []
        for name in [s for s in ["S1", "S2", "S3"] if s in scen]:
            med, _, _ = _mlu(scen[name]["cum_confirmed"])
            bits.append(f"{RESP.get(name, name)} {float(med[-1]):,.0f}")
        rows.append(("Step 3 — Forecast (cum. confirmed)", "; ".join(bits)))
    probs = _g(ss, "eoo_probs")
    plc = _g(ss, "eoo_valid_plc")
    thr = float(_g(ss, "eoo_threshold", 0.95))
    if probs is not None and plc:
        days = np.asarray(_g(ss, "eoo_days"))
        cross = np.where(np.asarray(probs) >= thr)[0]
        if len(cross) > 0:
            rows.append(("Step 4 — End of outbreak",
                         f"P>= {thr:.0%} at day {int(days[cross[0]])} after "
                         "last case"))
        else:
            rows.append(("Step 4 — End of outbreak",
                         f"P does not reach {thr:.0%} in window"))
    if not rows:
        rows = [("No results yet", "Run the steps in the app, then export.")]
    _table(slide, rows, Inches(0.6), Inches(1.4), Inches(12.1))
    _footer(slide, 2)


def _methods_slide(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "Methods & references")
    refs = [
        "Cori et al. 2013 — instantaneous R_t (Am J Epidemiol 178:1505).",
        "Nouvellet et al. 2018 — renewal-equation projection (Epidemics 22:3).",
        "WHO Ebola Response Team 2014 — natural history (NEJM 371:1481).",
        "Lloyd-Smith et al. 2005 — offspring dispersion k (Nature 438:355).",
        "Wamala et al. 2010 — onset-to-death lag (Emerg Infect Dis 16:1087).",
        "Nishiura et al. — end-of-outbreak probability framework.",
    ]
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12), Inches(4.5))
    tf = tb.text_frame; tf.word_wrap = True
    for i, ref in enumerate(refs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = "•  " + ref
        r.font.size = Pt(13); r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(8)
    cav = slide.shapes.add_textbox(Inches(0.6), Inches(6.0), Inches(12), Inches(1))
    cp = cav.text_frame.paragraphs[0]
    cr = cp.add_run()
    cr.text = ("Caveats: interpolated daily counts; R_t scenarios are "
               "assumptions, not predictions; PI bands reflect R_t sampling "
               "only.")
    cr.font.size = Pt(10); cr.font.italic = True; cr.font.color.rgb = GREY
    _footer(slide, page)


# ==========================================================================
# Public entry point
# ==========================================================================
def build_pptx(ss, scenario_name: str = "EVD Forecaster run") -> bytes:
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    now = datetime.now().strftime("%d %b %Y %H:%M")

    _title_slide(prs, scenario_name, now)
    _summary_slide(prs, ss, scenario_name)
    page = 3

    series = _g(ss, "result_series")
    snaps = _g(ss, "result_chart_snaps")
    if series is not None and len(series) > 0:
        _step_slide(prs, "Step 1 — Daily incidence", step1_inputs(ss),
                    render_daily(series), step1_interp(series, ss), page)
        page += 1
        _step_slide(prs, "Step 1 — Cumulative trend", step1_inputs(ss),
                    render_cumulative(series, snaps),
                    step1_interp(series, ss), page)
        page += 1

    rt_df = _g(ss, "rt_df")
    if rt_df is not None and len(rt_df) > 0:
        si = float(_g(ss, "rt_si_mean", 15.3))
        _step_slide(prs, "Step 2 — R_t estimation", step2_inputs(ss),
                    render_rt(rt_df, si), step2_interp(rt_df, si, ss), page)
        page += 1

    scen = _g(ss, "fc_scenarios")
    dates = _g(ss, "fc_dates")
    if scen and dates is not None:
        baselines = _g(ss, "fc_baselines", {})
        _step_slide(prs, "Step 3 — Forecast", step3_inputs(ss),
                    render_forecast(scen, dates, baselines, RESP, y_log=True),
                    step3_interp(scen, dates, ss), page)
        page += 1

    probs = _g(ss, "eoo_probs")
    plc = _g(ss, "eoo_valid_plc")
    if probs is not None and plc:
        thr = float(_g(ss, "eoo_threshold", 0.95))
        days = _g(ss, "eoo_days")
        _step_slide(prs, "Step 4 — End of outbreak", step4_inputs(ss),
                    render_eoo(days, probs, plc, thr),
                    step4_interp(days, probs, plc, thr), page)
        page += 1

    _methods_slide(prs, page)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
